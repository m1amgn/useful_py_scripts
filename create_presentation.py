import collections 
import collections.abc

from pptx import Presentation
from pptx.util import Inches

# Slide content
slide_content = [
    ("Introduction", "Phosphorus in Plant Nutrition\n\n- Essential nutrient for plant growth and development.\n- Soil phosphorus content impacts plant yield and quality of agricultural products."),
    ("Phosphorus Uptake by Plants", "Amount of phosphorus consumed by plants during vegetation.\nExample: Wheat consumes 10-12 kg of phosphorus per 1 ton of grain."),
    ("Phosphorus Fertilizers", "Role of Phosphate Fertilizers\n\n- To replenish phosphorus in soil for plant growth.\n- Necessary due to rapid binding of phosphorus in soil."),
]

# Create a PowerPoint presentation
prs = Presentation()

# Loop through the slide content and create slides
for title, content in slide_content:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title

    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(5.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    p = text_frame.add_paragraph()
    p.text = content

# Save the presentation
prs.save("presentation.pptx")
print("Presentation created and saved as 'presentation.pptx'.")
